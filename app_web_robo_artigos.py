from __future__ import annotations

from io import BytesIO
from html import escape
import re
import subprocess
import sys
import tempfile
import time
import unicodedata
from datetime import datetime
from pathlib import Path

import requests
import streamlit as st
import streamlit.components.v1 as components
from openpyxl import load_workbook
from pptx import Presentation


BASE = Path(__file__).resolve().parent
SEARCH_SCRIPT = BASE / "buscar_artigos_redes.py"
CLEAN_SCRIPT = BASE / "rpa_limpar_abstracts.py"
PRISMA_TEMPLATE = BASE / "prisma_template.pptx"
DEFAULT_PER_QUERY = 25
DEFAULT_WORKERS = 8
DEFAULT_MAX_RECORDS = 500
DEFAULT_TIMEOUT_MINUTES = 20
DEFAULT_INCLUDE_NETWORKS = False
TRANSLATION_GLOSSARY = {
    "aprendizado de maquina": "machine learning",
    "aprendizagem de maquina": "machine learning",
    "inteligencia artificial": "artificial intelligence",
    "redes neurais": "neural networks",
    "rede neural": "neural network",
    "redes neurais profundas": "deep neural networks",
    "aprendizado profundo": "deep learning",
    "aprendizagem profunda": "deep learning",
    "varejo": "retail",
    "varejo digital": "digital retail",
    "omnichannel": "omnichannel",
    "omnichannel varejo": "omnichannel retail",
    "cadeia de suprimentos": "supply chain",
    "logistica": "logistics",
    "comportamento do consumidor": "consumer behavior",
    "comportamento de compra": "shopping behavior",
    "experiencia do cliente": "customer experience",
    "satisfacao do cliente": "customer satisfaction",
    "fidelizacao de clientes": "customer loyalty",
    "transformacao digital": "digital transformation",
    "inovacao": "innovation",
    "sustentabilidade": "sustainability",
    "marketing digital": "digital marketing",
    "comercio eletronico": "e-commerce",
    "ecommerce": "e-commerce",
    "plataformas digitais": "digital platforms",
    "analise de redes": "network analysis",
    "redes sociais": "social networks",
    "grafos": "graphs",
    "redes bayesianas": "bayesian networks",
}
PRISMA_FIELDS = {
    "identified": "Registros identificados",
    "keyword_filter": "Registros apos filtro de palavras-chave",
    "screened": "Registros triados",
    "topic_not_addressed": "Tema nao abordado",
    "non_relevant": "Artigos nao relevantes",
    "retracted": "Artigos retratados",
    "abstract_screened": "Abstracts triados",
    "abstract_existent": "Abstracts existentes",
    "abstract_collected": "Abstracts coletados",
    "abstract_not_addressed": "Sem abstract valido",
    "author_screened": "Autores triados",
    "authors_standardized": "Autores padronizados",
    "records_kept": "Registros mantidos",
    "records_excluded": "Registros excluidos",
    "included": "Estudos incluidos",
}


def safe_filename(text: str) -> str:
    text = re.sub(r"[^a-zA-Z0-9_-]+", "_", text.strip())
    text = re.sub(r"_+", "_", text).strip("_")
    return text[:80] or f"planilha_artigos_{datetime.now():%Y%m%d}"


def normalize_lookup(text: str) -> str:
    text = unicodedata.normalize("NFKD", text.lower())
    text = "".join(ch for ch in text if not unicodedata.combining(ch))
    return re.sub(r"\s+", " ", text).strip()


def translate_pt_to_en(text: str) -> tuple[str, str]:
    text = text.strip()
    if not text:
        return "", ""

    lookup = normalize_lookup(text)
    if lookup in TRANSLATION_GLOSSARY:
        return TRANSLATION_GLOSSARY[lookup], "glossario"

    try:
        response = requests.get(
            "https://api.mymemory.translated.net/get",
            params={"q": text, "langpair": "pt|en"},
            timeout=8,
        )
        if response.status_code == 200:
            data = response.json()
            translated = (data.get("responseData") or {}).get("translatedText", "")
            translated = re.sub(r"\s+", " ", translated).strip()
            if translated:
                return translated, "tradutor online"
    except Exception:
        pass

    words = [TRANSLATION_GLOSSARY.get(normalize_lookup(part), part) for part in re.split(r"[,;]", text)]
    translated = ", ".join(part.strip() for part in words if part.strip())
    return translated or text, "texto original"


def command_text(cmd: list[str]) -> str:
    return " ".join(f'"{part}"' if " " in part else part for part in cmd)


def run_command(cmd: list[str], timeout_seconds: int) -> tuple[int, str]:
    started = time.monotonic()
    lines: list[str] = [f"$ {command_text(cmd)}\n"]
    process = subprocess.Popen(
        cmd,
        cwd=str(BASE),
        stdout=subprocess.PIPE,
        stderr=subprocess.STDOUT,
        text=True,
        encoding="utf-8",
        errors="replace",
    )
    assert process.stdout is not None

    try:
        while True:
            line = process.stdout.readline()
            if line:
                lines.append(line)
            if process.poll() is not None:
                remainder = process.stdout.read()
                if remainder:
                    lines.append(remainder)
                return process.returncode or 0, "".join(lines)
            if time.monotonic() - started > timeout_seconds:
                process.terminate()
                lines.append(f"\n[tempo limite atingido: {timeout_seconds // 60} min]\n")
                return 124, "".join(lines)
            time.sleep(0.1)
    finally:
        if process.poll() is None:
            process.kill()


def build_commands(
    keywords: str,
    output_dir: Path,
    filename_base: str,
    per_query: int,
    workers: int,
    max_records: int,
    year_from: int | None,
    year_to: int | None,
    include_networks: bool,
) -> tuple[Path, list[list[str]]]:
    raw_output = output_dir / f"{filename_base}_bruto.xlsx"
    final_output = output_dir / f"{filename_base}.xlsx"

    search_cmd = [
        sys.executable,
        "-u",
        str(SEARCH_SCRIPT),
        "--mode",
        "both" if include_networks else "keywords",
        "--keywords",
        keywords,
        "--per-query",
        str(per_query),
        "--workers",
        str(workers),
        "--max-records",
        str(max_records),
        "--wait",
        "0.05",
        "--output",
        str(raw_output),
    ]
    if year_from:
        search_cmd.extend(["--year-from", str(year_from)])
    if year_to:
        search_cmd.extend(["--year-to", str(year_to)])

    clean_cmd = [
        sys.executable,
        "-u",
        str(CLEAN_SCRIPT),
        "--input",
        str(raw_output),
        "--output",
        str(final_output),
        "--wait",
        "0.05",
        "--no-auto-reference",
        "--drop-empty-abstracts",
        "--skip-missing-fetch",
    ]
    return final_output, [search_cmd, clean_cmd]


def read_final_count(workbook_path: Path) -> int:
    try:
        wb = load_workbook(workbook_path, read_only=True, data_only=True)
        ws = wb.active
        return max(ws.max_row - 1, 0)
    except Exception:
        return 0


def extract_int_from_logs(logs: list[str], label: str) -> int:
    joined = "\n".join(logs)
    match = re.search(rf"{re.escape(label)}:\s*(\d+)", joined)
    return int(match.group(1)) if match else 0


def build_prisma_values(logs: list[str], workbook_path: Path) -> dict[str, int]:
    included = read_final_count(workbook_path)
    discarded_without_abstract = extract_int_from_logs(logs, "Artigos sem abstract descartados")
    identified = max(included + discarded_without_abstract, included)
    return {
        "identified": identified,
        "keyword_filter": identified,
        "screened": identified,
        "topic_not_addressed": discarded_without_abstract,
        "non_relevant": 0,
        "retracted": 0,
        "abstract_screened": identified,
        "abstract_existent": included,
        "abstract_collected": 0,
        "abstract_not_addressed": discarded_without_abstract,
        "author_screened": included,
        "authors_standardized": included,
        "records_kept": included,
        "records_excluded": 0,
        "included": included,
    }


def prisma_texts(values: dict[str, int]) -> dict[str, str]:
    return {
        "heading": "Identification of studies via bibliographic data search engine",
        "identified": f"Records identified from academic databases\n(n={values['identified']})",
        "keyword_filter": f"Keyword search terms matched\n(n={values['keyword_filter']})",
        "screened": f"Records screened\n(n={values['screened']})",
        "exclusion": (
            f"Research topic not addressed (n={values['topic_not_addressed']}); "
            f"Non-relevant articles (n={values['non_relevant']}); "
            f"Retracted articles (n={values['retracted']})"
        ),
        "abstract_screened": f"Records (Abstracts) screened\n(n={values['abstract_screened']})",
        "abstract_status": (
            f"Existent (n={values['abstract_existent']}); "
            f"Collected (n={values['abstract_collected']}); "
            f"Not addressed (n={values['abstract_not_addressed']})"
        ),
        "author_screened": f"Records (Authors) screened\n(n={values['author_screened']})",
        "author_status": (
            f"Standardization (given and family names) (n={values['authors_standardized']}); "
            f"Records kept (n={values['records_kept']}); "
            f"Records excluded (n={values['records_excluded']})"
        ),
        "included": f"Studies included in review\n(n={values['included']})",
    }


def wrap_svg_text(text: str, x: int, y: int, width: int, line_height: int = 18, size: int = 15) -> str:
    words = text.replace("\n", " \n ").split()
    lines: list[str] = []
    current = ""
    max_chars = max(width // 8, 18)
    for word in words:
        if word == "\n":
            if current:
                lines.append(current)
                current = ""
            continue
        candidate = f"{current} {word}".strip()
        if len(candidate) > max_chars and current:
            lines.append(current)
            current = word
        else:
            current = candidate
    if current:
        lines.append(current)

    tspans = []
    for index, line in enumerate(lines[:5]):
        tspans.append(
            f'<tspan x="{x}" dy="{0 if index == 0 else line_height}">{escape(line)}</tspan>'
        )
    return f'<text x="{x}" y="{y}" font-size="{size}" fill="#0f172a" font-family="Arial, sans-serif">{"".join(tspans)}</text>'


def render_prisma_svg(values: dict[str, int]) -> str:
    texts = prisma_texts(values)
    boxes = [
        ("Identification", 70, 116, 130, 70, "#ecfeff"),
        ("Screening", 70, 332, 130, 70, "#eff6ff"),
        ("Inclusion", 70, 548, 130, 70, "#fefce8"),
    ]
    flow_boxes = [
        ("identified", 260, 112, 300, 90),
        ("keyword_filter", 720, 112, 330, 90),
        ("screened", 280, 258, 270, 64),
        ("exclusion", 720, 246, 330, 90),
        ("abstract_screened", 280, 374, 270, 70),
        ("abstract_status", 720, 366, 330, 78),
        ("author_screened", 280, 488, 270, 70),
        ("author_status", 720, 474, 330, 92),
        ("included", 270, 604, 290, 64),
    ]
    svg_parts = [
        '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 1120 720" width="1120" height="720">',
        '<rect width="1120" height="720" fill="#ffffff"/>',
        '<rect x="230" y="34" width="650" height="44" rx="10" fill="#dbeafe"/>',
        wrap_svg_text(texts["heading"], 260, 62, 600, size=17),
    ]
    for label, x, y, w, h, fill in boxes:
        svg_parts.extend(
            [
                f'<rect x="{x}" y="{y}" width="{w}" height="{h}" rx="12" fill="{fill}" stroke="#94a3b8"/>',
                wrap_svg_text(label, x + 20, y + 42, w - 30, size=16),
            ]
        )
    for key, x, y, w, h in flow_boxes:
        svg_parts.extend(
            [
                f'<rect x="{x}" y="{y}" width="{w}" height="{h}" rx="12" fill="#ffffff" stroke="#2563eb" stroke-width="2"/>',
                wrap_svg_text(texts[key], x + 18, y + 30, w - 36),
            ]
        )
    arrows = [
        (560, 157, 720, 157),
        (410, 202, 410, 258),
        (550, 290, 720, 290),
        (410, 322, 410, 374),
        (550, 409, 720, 409),
        (410, 444, 410, 488),
        (550, 522, 720, 522),
        (410, 558, 410, 604),
    ]
    svg_parts.append(
        '<defs><marker id="arrow" markerWidth="10" markerHeight="10" refX="9" refY="3" orient="auto" markerUnits="strokeWidth"><path d="M0,0 L0,6 L9,3 z" fill="#0f766e"/></marker></defs>'
    )
    for x1, y1, x2, y2 in arrows:
        svg_parts.append(
            f'<line x1="{x1}" y1="{y1}" x2="{x2}" y2="{y2}" stroke="#0f766e" stroke-width="3" marker-end="url(#arrow)"/>'
        )
    svg_parts.append("</svg>")
    return "".join(svg_parts)


def generate_prisma_pptx(values: dict[str, int]) -> bytes:
    texts = prisma_texts(values)
    replacements = {
        "CaixaDeTexto 8": texts["heading"],
        "CaixaDeTexto 2": texts["identified"],
        "CaixaDeTexto 17": texts["keyword_filter"],
        "CaixaDeTexto 12": texts["screened"],
        "CaixaDeTexto 13": texts["exclusion"],
        "CaixaDeTexto 10": texts["abstract_screened"],
        "CaixaDeTexto 11": texts["abstract_status"],
        "CaixaDeTexto 14": texts["author_screened"],
        "CaixaDeTexto 15": texts["author_status"],
        "CaixaDeTexto 16": texts["included"],
    }
    prs = Presentation(PRISMA_TEMPLATE)
    slide = prs.slides[0]
    for shape in slide.shapes:
        if shape.name in replacements and hasattr(shape, "text_frame"):
            shape.text_frame.clear()
            paragraph = shape.text_frame.paragraphs[0]
            paragraph.text = replacements[shape.name]
    buffer = BytesIO()
    prs.save(buffer)
    return buffer.getvalue()


def render_prisma_section() -> None:
    if "prisma_values" not in st.session_state:
        return

    st.divider()
    st.subheader("PRISMA editavel")
    st.caption("Os campos abaixo foram preenchidos com base na planilha gerada. Voce pode ajustar os numeros antes de baixar.")

    values = st.session_state["prisma_values"].copy()
    cols = st.columns(3)
    updated: dict[str, int] = {}
    for index, (key, label) in enumerate(PRISMA_FIELDS.items()):
        with cols[index % 3]:
            updated[key] = int(
                st.number_input(label, min_value=0, value=int(values.get(key, 0)), step=1, key=f"prisma_{key}")
            )
    st.session_state["prisma_values"] = updated

    svg = render_prisma_svg(updated)
    components.html(svg, height=640, scrolling=False)
    pptx_bytes = generate_prisma_pptx(updated)
    st.download_button(
        "Baixar PRISMA em PowerPoint editavel",
        data=pptx_bytes,
        file_name="prisma_fluxo_revisao.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )
    st.download_button(
        "Baixar PRISMA em SVG vetorial",
        data=svg.encode("utf-8"),
        file_name="prisma_fluxo_revisao.svg",
        mime="image/svg+xml",
    )


def apply_theme() -> None:
    st.markdown(
        """
        <style>
        .stApp {
            background: linear-gradient(180deg, #f6fbff 0%, #ffffff 46%, #f8fafc 100%);
        }
        .block-container {
            padding-top: 2rem;
            padding-bottom: 3rem;
            max-width: 980px;
        }
        .hero {
            border-radius: 18px;
            padding: 30px 34px;
            margin-bottom: 24px;
            background:
                linear-gradient(135deg, rgba(13, 148, 136, 0.94), rgba(37, 99, 235, 0.92) 58%, rgba(245, 158, 11, 0.94));
            color: white;
            box-shadow: 0 18px 40px rgba(15, 23, 42, 0.16);
        }
        .hero h1 {
            color: white;
            font-size: 2.2rem;
            line-height: 1.1;
            margin: 0 0 10px;
            letter-spacing: 0;
        }
        .hero p {
            color: rgba(255, 255, 255, 0.94);
            font-size: 1.02rem;
            margin: 0;
            max-width: 760px;
        }
        .eyebrow {
            display: inline-block;
            padding: 5px 10px;
            border-radius: 999px;
            background: rgba(255, 255, 255, 0.18);
            color: white;
            font-weight: 700;
            font-size: 0.78rem;
            margin-bottom: 14px;
        }
        .step-grid {
            display: grid;
            grid-template-columns: repeat(3, minmax(0, 1fr));
            gap: 12px;
            margin: 0 0 22px;
        }
        .step-card {
            border: 1px solid #dbeafe;
            border-radius: 12px;
            padding: 14px 16px;
            background: white;
            box-shadow: 0 8px 24px rgba(15, 23, 42, 0.06);
        }
        .step-number {
            color: #0f766e;
            font-weight: 800;
            font-size: 0.82rem;
            margin-bottom: 4px;
        }
        .step-card strong {
            display: block;
            color: #0f172a;
            margin-bottom: 3px;
        }
        .step-card span {
            color: #475569;
            font-size: 0.9rem;
        }
        .section-copy {
            color: #475569;
            margin: -6px 0 14px;
        }
        .stTextInput input {
            border-radius: 10px;
            border: 1px solid #cbd5e1;
            background: #ffffff;
        }
        .stTextInput input:focus {
            border-color: #0d9488;
            box-shadow: 0 0 0 1px #0d9488;
        }
        .stButton button, .stDownloadButton button {
            border-radius: 999px;
            border: 0;
            background: linear-gradient(90deg, #0f766e, #2563eb);
            color: white;
            font-weight: 800;
            padding: 0.7rem 1.25rem;
        }
        .stButton button:hover, .stDownloadButton button:hover {
            color: white;
            filter: brightness(1.05);
        }
        [data-testid="stFormSubmitButton"] button {
            width: 100%;
            min-height: 56px;
            margin-top: 10px;
            font-size: 1.05rem;
            letter-spacing: 0;
            background: linear-gradient(90deg, #f59e0b, #0d9488 45%, #2563eb);
            box-shadow: 0 14px 30px rgba(37, 99, 235, 0.22);
        }
        [data-testid="stFormSubmitButton"] button:hover {
            box-shadow: 0 16px 36px rgba(13, 148, 136, 0.28);
            transform: translateY(-1px);
        }
        [data-testid="stForm"] {
            border: 1px solid #dbeafe;
            border-radius: 16px;
            padding: 22px;
            background: rgba(255, 255, 255, 0.92);
            box-shadow: 0 10px 30px rgba(15, 23, 42, 0.06);
        }
        @media (max-width: 760px) {
            .step-grid {
                grid-template-columns: 1fr;
            }
            .hero {
                padding: 24px;
            }
        }
        </style>
        """,
        unsafe_allow_html=True,
    )


def main() -> None:
    st.set_page_config(page_title="Robo de artigos academicos", page_icon=None, layout="centered")
    apply_theme()
    for index in range(1, 6):
        st.session_state.setdefault(f"keyword_{index}", "")

    st.markdown(
        """
        <div class="hero">
            <div class="eyebrow">Busca academica automatizada</div>
            <h1>Gerador de planilha de artigos academicos</h1>
            <p>Informe seus termos de pesquisa. O aplicativo busca artigos em bases abertas, limpa os abstracts e entrega um Excel padronizado para baixar.</p>
        </div>
        """,
        unsafe_allow_html=True,
    )
    st.markdown(
        """
        <div class="step-grid">
            <div class="step-card"><div class="step-number">Etapa 1</div><strong>Busca</strong><span>Consulta bases academicas abertas.</span></div>
            <div class="step-card"><div class="step-number">Etapa 2</div><strong>Organiza</strong><span>Remove duplicados e abstracts invalidos.</span></div>
            <div class="step-card"><div class="step-number">Etapa 3</div><strong>Entrega</strong><span>Gera uma planilha Excel pronta.</span></div>
        </div>
        """,
        unsafe_allow_html=True,
    )

    input_col, translator_col = st.columns([1.55, 1], gap="large")

    with input_col:
        with st.form("search_form"):
            st.subheader("Tema ou palavras-chave")
            st.markdown('<p class="section-copy">Insira ate 5 palavras-chave ou consultas. Cada campo preenchido gera uma busca separada.</p>', unsafe_allow_html=True)
            for index in range(1, 6):
                st.text_input(
                    f"Palavra-chave {index}",
                    key=f"keyword_{index}",
                    placeholder="Ex.: artificial intelligence in education",
                )

            c1, c2 = st.columns(2)
            year_from_raw = c1.text_input("Ano inicial opcional", value="", placeholder="Ex.: 2020")
            year_to_raw = c2.text_input("Ano final opcional", value="", placeholder="Ex.: 2026")

            filename_base = st.text_input(
                "Nome do arquivo",
                value="planilha_artigos_academicos",
                placeholder="Ex.: revisao_inteligencia_artificial",
                help="Voce pode manter o nome sugerido ou escrever outro. A extensao .xlsx sera adicionada automaticamente.",
            )
            st.info(
                "Ao clicar, o aplicativo vai consultar fontes academicas abertas, combinar resultados repetidos, "
                "limpar abstracts e montar uma planilha Excel padronizada para download."
            )
            submitted = st.form_submit_button("Gerar minha planilha Excel pronta")

    with translator_col:
        with st.container(border=True):
            st.subheader("Tradutor rapido")
            st.caption("Digite uma palavra ou expressao em portugues para usar na busca em ingles.")
            pt_text = st.text_input("Termo em portugues", placeholder="Ex.: varejo digital")
            if st.button("Traduzir para ingles", use_container_width=True):
                translated, source = translate_pt_to_en(pt_text)
                st.session_state["translated_keyword"] = translated
                st.session_state["translated_source"] = source

            translated_keyword = st.session_state.get("translated_keyword", "")
            if translated_keyword:
                st.success(translated_keyword)
                st.caption(f"Fonte: {st.session_state.get('translated_source', 'tradutor')}")
                if st.button("Inserir no proximo campo livre", use_container_width=True):
                    inserted = False
                    for index in range(1, 6):
                        key = f"keyword_{index}"
                        if not st.session_state.get(key, "").strip():
                            st.session_state[key] = translated_keyword
                            inserted = True
                            break
                    if not inserted:
                        st.session_state["keyword_5"] = translated_keyword
                    st.rerun()

    if "excel_bytes" in st.session_state and not submitted:
        st.download_button(
            "Baixar ultimo Excel gerado",
            data=st.session_state["excel_bytes"],
            file_name=st.session_state.get("excel_name", "planilha_artigos_academicos.xlsx"),
            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        )
        render_prisma_section()

    if not submitted:
        return

    keywords = "\n".join(st.session_state.get(f"keyword_{index}", "").strip() for index in range(1, 6) if st.session_state.get(f"keyword_{index}", "").strip())
    if not keywords.strip():
        st.error("Digite pelo menos um tema ou palavra-chave.")
        return

    try:
        year_from = int(year_from_raw) if year_from_raw.strip() else None
        year_to = int(year_to_raw) if year_to_raw.strip() else None
    except ValueError:
        st.error("Ano inicial e ano final precisam ser numeros inteiros.")
        return

    safe_base = safe_filename(filename_base)
    timeout_seconds = DEFAULT_TIMEOUT_MINUTES * 60

    with tempfile.TemporaryDirectory(prefix="robo_artigos_") as tmp:
        output_dir = Path(tmp)
        final_output, commands = build_commands(
            keywords=keywords,
            output_dir=output_dir,
            filename_base=safe_base,
            per_query=DEFAULT_PER_QUERY,
            workers=DEFAULT_WORKERS,
            max_records=DEFAULT_MAX_RECORDS,
            year_from=year_from,
            year_to=year_to,
            include_networks=DEFAULT_INCLUDE_NETWORKS,
        )

        all_logs: list[str] = []
        with st.status("Preparando sua planilha...", expanded=False) as status:
            for index, cmd in enumerate(commands, start=1):
                if index == 1:
                    st.write("Buscando artigos e abstracts nas bases academicas.")
                else:
                    st.write("Limpando e organizando a planilha final.")
                code, log = run_command(cmd, timeout_seconds)
                all_logs.append(log)
                if code != 0:
                    status.update(label="Execucao interrompida", state="error")
                    st.error("A planilha nao foi concluida. Tente uma busca mais especifica ou execute novamente.")
                    st.session_state["last_log"] = "\n".join(all_logs)
                    return
            status.update(label="Planilha pronta", state="complete")

        if not final_output.exists():
            st.error("O arquivo final nao foi encontrado apos a execucao.")
            return

        excel_bytes = final_output.read_bytes()
        st.session_state["excel_bytes"] = excel_bytes
        st.session_state["excel_name"] = f"{safe_base}.xlsx"
        st.session_state["last_log"] = "\n".join(all_logs)
        st.session_state["prisma_values"] = build_prisma_values(all_logs, final_output)

        st.success("Sua planilha esta pronta.")
        st.download_button(
            "Baixar planilha Excel",
            data=excel_bytes,
            file_name=f"{safe_base}.xlsx",
            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        )
        render_prisma_section()


if __name__ == "__main__":
    main()
